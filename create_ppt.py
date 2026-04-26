from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def create_presentation():
    # Create a presentation object
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    # ------------------ Slide 1: Title ------------------
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "Rainfall Timeseries Forecasting"
    subtitle.text = "Monthly Rainfall Forecast (Mar-Oct 2026)\n& Daily Probability Interface"
    
    # ------------------ Slide 2: Project Objectives & Highlights ------------------
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide2.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Project Objectives & Highlights"
    
    tf = body_shape.text_frame
    tf.text = "Compare multiple forecasting models with a time-aware validation setup"
    
    p = tf.add_paragraph()
    p.text = "Select the final model using strict validation metrics (avoids test leakage)"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Persist robust outputs: forecasts, evaluation tables, figures, and metadata"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Serve results through a clean, interactive Streamlit application"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Keep heavy training dependencies strictly separated from the lightweight app"
    p.level = 0

    # ------------------ Slide 3: Best Model Performance ------------------
    slide3 = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide3.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Best Saved Model"
    
    tf = body_shape.text_frame
    tf.text = "Selected Model Architecture: BiLSTM"
    
    p = tf.add_paragraph()
    p.text = "Validation RMSE: 88.26"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Test RMSE: 110.92"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Forecast Output: Future forecasts generated and saved successfully."
    p.level = 0
    
    # ------------------ Slide 4: Project Methodology & Pipeline ------------------
    slide4 = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide4.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Methodology & Pipeline"
    
    tf = body_shape.text_frame
    tf.text = "Data Processing:"
    
    p = tf.add_paragraph()
    p.text = "Extracts daily behavior and creates a monthly dataset with climatology profiles"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Model Training:"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Trains on historical datasets with robust sequence models"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Evaluation:"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Produces horizon-wise evaluation metrics and backtest predictions"
    p.level = 1
    
    # ------------------ Slide 5: The Streamlit App ------------------
    slide5 = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide5.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Interactive Application"
    
    tf = body_shape.text_frame
    tf.text = "Local App Entry: Runs via app.py using Streamlit"
    
    p = tf.add_paragraph()
    p.text = "Key Features:"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Visualizes the 2026 future forecast"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Interactive exploration of evaluation figures and tables"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Estimates which days in a selected month have the highest chance of rainfall based on historical data"
    p.level = 1
    
    # ------------------ Slide 6: Conclusion ------------------
    slide6 = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide6.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Conclusion & Next Steps"
    
    tf = body_shape.text_frame
    tf.text = "Delivers accurate, backtested rainfall forecasting using modern Deep Learning (BiLSTM)"
    
    p = tf.add_paragraph()
    p.text = "User-friendly and accessible through a scalable web interface"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Modular codebase is easy to extend with new models, pipelines, or data updates"
    p.level = 0

    # Save presentation
    output_path = os.path.join(os.getcwd(), 'Project_Presentation.pptx')
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == '__main__':
    create_presentation()
